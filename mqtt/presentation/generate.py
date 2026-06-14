import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import os

OUT = os.path.dirname(__file__)

# ── Colors ──────────────────────────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT = RGBColor(0x1A, 0x56, 0xDB)   # blue

CHART_COLORS = ['#1A56DB', '#4B9CD3', '#A8C8F0']

def set_bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_title(slide, text, top=Inches(0.35), size=28):
    txb = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.6))
    tf  = txb.text_frame
    tf.word_wrap = False
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = BLACK

def add_text(slide, text, left, top, width, height, size=11, bold=False, align=PP_ALIGN.LEFT, color=BLACK):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def chart_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf

def style_ax(ax, title=''):
    ax.set_facecolor('white')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#CCCCCC')
    ax.spines['bottom'].set_color('#CCCCCC')
    ax.tick_params(colors='#333333', labelsize=9)
    ax.yaxis.label.set_color('#333333')
    ax.xaxis.label.set_color('#333333')
    if title:
        ax.set_title(title, fontsize=10, fontweight='bold', color='#111111', pad=8)

# ── Data ────────────────────────────────────────────────────────────────
# Scenario A
sa_qos   = [0,0,0,1,1,1,2,2,2]
sa_dev   = [100,1000,10000]*3
sa_tput  = [100,1000,10000, 100,1000,10000, 100,1000,11000]
sa_cpu   = [2.62,2.70,12.61, 5.24,4.79,29.77, 2.70,6.26,102.83]

# Scenario C
sc_phases  = ['Normal\n(50 dev)','Burst\n(5000 dev)','Recovery\n(50 dev)']
sc_tput_q0 = [55, 5000, 50]
sc_tput_q1 = [55, 5000, 50]
sc_cpu_q0  = [1.65, 9.43, 1.93]
sc_cpu_q1  = [1.98, 18.09, 2.02]

# Scenario D
sd_qos_labels = ['QoS 0','QoS 1','QoS 2']
sd_p50  = [3.9, 5.0, 4.1]
sd_p95  = [13.2, 28.4, 31.4]
sd_p99  = [13.9, 33.4, 35.9]

# Benchmark
bm_clients = [100, 500, 1000]
bm_q0_avg  = [197774, 623658, 576370]
bm_q1_avg  = [79816,  74638,  71862]
bm_q2_avg  = [42919,  40604,  38983]
bm_q0_cpu  = [62.89, 99.94, 99.51]
bm_q1_cpu  = [106.41,105.83,104.54]
bm_q2_cpu  = [105.75,105.91,104.80]

# ── Presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

blank = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)

add_text(sl, 'IoT Mikroservisi zasnovani na dogadjajima',
         Inches(1), Inches(1.6), Inches(8), Inches(0.8),
         size=30, bold=True, align=PP_ALIGN.CENTER)

add_text(sl, 'MQTT vs Kafka — Eksperimentalna analiza performansi',
         Inches(1), Inches(2.5), Inches(8), Inches(0.5),
         size=16, align=PP_ALIGN.CENTER, color=RGBColor(0x44,0x44,0x44))

add_text(sl, 'Spring Boot (MQTT)  |  .NET Core (Kafka)  |  PostgreSQL',
         Inches(1), Inches(3.2), Inches(8), Inches(0.4),
         size=12, align=PP_ALIGN.CENTER, color=RGBColor(0x77,0x77,0x77))

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Arhitektura
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Arhitektura sistema')

arch_text = (
    "MQTT strana (Spring Boot)\n"
    "  Ingestion Service  →  Mosquitto Broker  →  Storage Service  →  PostgreSQL\n"
    "                                          →  Analytics Service (tumbling window 10s)\n\n"
    "Kafka strana (.NET Core)\n"
    "  Ingestion Service  →  Kafka KRaft      →  Storage Service  →  PostgreSQL\n"
    "                                          →  Analytics Service\n\n"
    "Zajednička baza: PostgreSQL (schema: iot_mqtt / iot_kafka)\n"
    "Dataset: real_time_data.csv  |  10 atributa  |  Alarm: temperatura > 50 C"
)

txb = sl.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4))
tf  = txb.text_frame
tf.word_wrap = True
for i, line in enumerate(arch_text.split('\n')):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(11)
    run.font.color.rgb = BLACK
    if line and not line.startswith(' '):
        run.font.bold = True

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Scenario A: Throughput
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Scenario A — Masovni unos podataka (Throughput)')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.6))
fig.patch.set_facecolor('white')

# Chart 1: throughput by device count per QoS
x = np.arange(3)
w = 0.25
labels = ['100', '1 000', '10 000']
for i, (qos, color) in enumerate(zip([0,1,2], CHART_COLORS)):
    vals = [sa_tput[i*3], sa_tput[i*3+1], sa_tput[i*3+2]]
    axes[0].bar(x + i*w, vals, w, label=f'QoS {qos}', color=color)
axes[0].set_xticks(x + w)
axes[0].set_xticklabels(labels)
axes[0].set_xlabel('Broj uredjaja')
axes[0].set_ylabel('msg/s')
style_ax(axes[0], 'Throughput po broju uredjaja')
axes[0].legend(fontsize=8)

# Chart 2: CPU usage QoS 2 vs QoS 0 at 10000 devices
qos_labels = ['QoS 0', 'QoS 1', 'QoS 2']
cpu_10k = [12.61, 29.77, 102.83]
bars = axes[1].bar(qos_labels, cpu_10k, color=CHART_COLORS, width=0.5)
axes[1].set_ylabel('CPU %')
axes[1].axhline(100, color='#CC0000', linestyle='--', linewidth=1, label='100% (1 core)')
axes[1].legend(fontsize=8)
style_ax(axes[1], 'CPU na 10 000 uredjaja')
for bar, val in zip(bars, cpu_10k):
    axes[1].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                 f'{val}%', ha='center', va='bottom', fontsize=9)

plt.tight_layout(pad=1.0)
img = chart_to_image(fig)
sl.shapes.add_picture(img, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3))

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Scenario C: Burst Load
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Scenario C — Nagli porast opterecenja (Burst Load)')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.6))
fig.patch.set_facecolor('white')

x = np.arange(3)
w = 0.35
axes[0].bar(x - w/2, sc_tput_q0, w, label='QoS 0', color=CHART_COLORS[0])
axes[0].bar(x + w/2, sc_tput_q1, w, label='QoS 1', color=CHART_COLORS[1])
axes[0].set_xticks(x)
axes[0].set_xticklabels(sc_phases)
axes[0].set_ylabel('msg/s')
style_ax(axes[0], 'Throughput po fazi')
axes[0].legend(fontsize=8)

axes[1].bar(x - w/2, sc_cpu_q0, w, label='QoS 0', color=CHART_COLORS[0])
axes[1].bar(x + w/2, sc_cpu_q1, w, label='QoS 1', color=CHART_COLORS[1])
axes[1].set_xticks(x)
axes[1].set_xticklabels(sc_phases)
axes[1].set_ylabel('CPU %')
style_ax(axes[1], 'CPU po fazi')
axes[1].legend(fontsize=8)

plt.tight_layout(pad=1.0)
img = chart_to_image(fig)
sl.shapes.add_picture(img, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3))

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Scenario B: Network disconnect
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Scenario B — Pad mreze i automatski oporavak')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.6))
fig.patch.set_facecolor('white')

# Timeline: sentCount
phases_b   = ['Baseline', 'Disconnect\n(30s)', 'Recovery\n+10s', 'Recovery\n+20s']
sent_b     = [137600, None, 141900, 143000]
records_b  = [29499,  None, None,   30599]

x_b = [0, 1, 2, 3]
sent_vals  = [137600, 137600, 141900, 143000]
axes[0].plot(x_b, sent_vals, 'o-', color=CHART_COLORS[0], linewidth=2, markersize=6)
axes[0].axvspan(0.5, 1.5, color='#FFCCCC', alpha=0.5, label='Disconnect interval')
axes[0].set_xticks(x_b)
axes[0].set_xticklabels(phases_b, fontsize=8)
axes[0].set_ylabel('Ukupno poruka poslato')
style_ax(axes[0], 'sentCount tokom testa')
axes[0].legend(fontsize=8)
axes[0].yaxis.set_major_formatter(matplotlib.ticker.FuncFormatter(lambda x, _: f'{int(x):,}'))

# DB records delta
categories = ['Pre\ndisconnecta', 'Posle\nreconnecta\n(+20s)']
db_vals    = [29499, 30599]
bars = axes[1].bar(categories, db_vals, color=[CHART_COLORS[1], CHART_COLORS[0]], width=0.4)
axes[1].set_ylabel('Broj zapisa u bazi')
style_ax(axes[1], 'DB zapisi — 1 100 novih posle reconnecta')
for bar, val in zip(bars, db_vals):
    axes[1].text(bar.get_x() + bar.get_width()/2, bar.get_height() + 100,
                 f'{val:,}', ha='center', fontsize=9)
axes[1].set_ylim(0, 35000)

plt.tight_layout(pad=1.0)
img = chart_to_image(fig)
sl.shapes.add_picture(img, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3))

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Scenario D: Latencija
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Scenario D — End-to-end latencija (100 uzoraka po QoS)')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.6))
fig.patch.set_facecolor('white')

x = np.arange(3)
w = 0.25
axes[0].bar(x - w, sd_p50, w, label='p50', color=CHART_COLORS[0])
axes[0].bar(x,     sd_p95, w, label='p95', color=CHART_COLORS[1])
axes[0].bar(x + w, sd_p99, w, label='p99', color=CHART_COLORS[2])
axes[0].set_xticks(x)
axes[0].set_xticklabels(sd_qos_labels)
axes[0].set_ylabel('ms')
style_ax(axes[0], 'Prosecna latencija (avg od 100 uzoraka)')
axes[0].legend(fontsize=8)

# p99 comparison samo
axes[1].bar(sd_qos_labels, sd_p99, color=CHART_COLORS, width=0.5)
axes[1].set_ylabel('ms')
style_ax(axes[1], 'p99 latencija po QoS nivou')
for i, (label, val) in enumerate(zip(sd_qos_labels, sd_p99)):
    axes[1].text(i, val + 0.3, f'{val} ms', ha='center', fontsize=9)

plt.tight_layout(pad=1.0)
img = chart_to_image(fig)
sl.shapes.add_picture(img, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3))

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Benchmark: Broker kapacitet
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Faza 5 — Broker benchmark (emqtt-bench, 256B poruke)')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.6))
fig.patch.set_facecolor('white')

x = np.arange(3)
w = 0.25
xlabels = ['100', '500', '1 000']
axes[0].bar(x - w, [v/1000 for v in bm_q0_avg], w, label='QoS 0', color=CHART_COLORS[0])
axes[0].bar(x,     [v/1000 for v in bm_q1_avg], w, label='QoS 1', color=CHART_COLORS[1])
axes[0].bar(x + w, [v/1000 for v in bm_q2_avg], w, label='QoS 2', color=CHART_COLORS[2])
axes[0].set_xticks(x)
axes[0].set_xticklabels(xlabels)
axes[0].set_xlabel('Broj klijenata')
axes[0].set_ylabel('msg/s (hiljada)')
style_ax(axes[0], 'Throughput brokera')
axes[0].legend(fontsize=8)

axes[1].bar(x - w, bm_q0_cpu, w, label='QoS 0', color=CHART_COLORS[0])
axes[1].bar(x,     bm_q1_cpu, w, label='QoS 1', color=CHART_COLORS[1])
axes[1].bar(x + w, bm_q2_cpu, w, label='QoS 2', color=CHART_COLORS[2])
axes[1].set_xticks(x)
axes[1].set_xticklabels(xlabels)
axes[1].set_xlabel('Broj klijenata')
axes[1].set_ylabel('CPU % (Mosquitto)')
axes[1].axhline(100, color='#CC0000', linestyle='--', linewidth=1)
style_ax(axes[1], 'CPU Mosquitto brokera')
axes[1].legend(fontsize=8)

plt.tight_layout(pad=1.0)
img = chart_to_image(fig)
sl.shapes.add_picture(img, Inches(0.3), Inches(1.0), Inches(9.4), Inches(4.3))

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Uporedna tabela
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Uporedna tabela performansi')

headers = ['', 'Throughput (max)', 'p95 latencija', 'CPU (max load)', 'RAM footprint']
rows = [
    ['MQTT QoS 0', '623k msg/s', '13 ms', '99%', '~870 MB'],
    ['MQTT QoS 1', '80k msg/s',  '28 ms', '106%', '~880 MB'],
    ['MQTT QoS 2', '43k msg/s',  '31 ms', '106%', '~870 MB'],
    ['Kafka (acks=0)', 'kolega',  'kolega', 'kolega', '~512 MB+'],
    ['Kafka (acks=all)', 'kolega', 'kolega', 'kolega', '~512 MB+'],
]

col_w = [1.8, 1.7, 1.5, 1.5, 1.6]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

row_h = 0.48
top_start = 1.1

# Header
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    txb = sl.shapes.add_textbox(Inches(x), Inches(top_start), Inches(w), Inches(row_h))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = h
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = BLACK

# Rows
for i, row in enumerate(rows):
    y = top_start + row_h * (i + 1)
    bg_color = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else WHITE
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        shape = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(row_h))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = cell
        run.font.size = Pt(9)
        run.font.bold = (j == 0)
        run.font.color.rgb = BLACK

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — MQTT: prednosti i ogranicenja
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'MQTT — Prednosti na edge, ogranicenja za analitiku')

left_items = [
    'Prednosti na edge:',
    '  2-byte header overhead',
    '  Radi na TCP sa malim RAM footprintom',
    '  QoS garantuje isporuku na nestabilnoj mrezi',
    '  cleanSession=false: broker cuva poruke dok je klijent offline',
    '  Automatski reconnect bez gubitka pretplate',
    '  623k msg/s na QoS 0 (benchmark)',
]

right_items = [
    'Ogranicenja za istorijsku analitiku:',
    '  Broker ne cuva istoriju poruka',
    '  Nema replay mehanizma',
    '  Retained message: samo poslednja vrednost po topicu',
    '  Horizontalno skaliranje brokera je kompleksno',
    '  Nema built-in stream processing',
    '  Nema consumer group koncepta',
]

for i, item in enumerate(left_items):
    bold = (i == 0)
    add_text(sl, item, Inches(0.4), Inches(1.0 + i*0.55), Inches(4.5), Inches(0.5),
             size=10, bold=bold)

for i, item in enumerate(right_items):
    bold = (i == 0)
    add_text(sl, item, Inches(5.1), Inches(1.0 + i*0.55), Inches(4.5), Inches(0.5),
             size=10, bold=bold)

# Divider line
line = sl.shapes.add_connector(1, Inches(5.0), Inches(1.0), Inches(5.0), Inches(5.2))
line.line.color.rgb = GRAY
line.line.width = Pt(1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Kafka vs MQTT
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Kafka — Cloud dominacija i cena skalabilnosti')

items = [
    ('Kafka prednosti:', [
        'Distribuirani commit log — sve poruke sacuvane (konfigurisani retention)',
        'Consumer grupe — nezavisno citanje istog streama',
        'Horizontalno skaliranje: dodavanjem brokera i particija',
        'Replay: reprocessing istorijskih podataka',
        '3 particije po topicu — device_id kao partition key',
    ]),
    ('Cena skalabilnosti:', [
        'Minimalni footprint ~512 MB RAM samo za JVM',
        'KRaft mod eliminise Zookeeper overhead',
        'Visoka latencija zbog batch commit logike vs MQTT <5ms p50',
        'Nerealno pokretati na ARM edge uredjaijima sa 256 MB RAM',
        'Kompleksna konfiguracija vs MQTT 5-linijski mosquitto.conf',
    ]),
]

y = 1.05
for title, subitems in items:
    add_text(sl, title, Inches(0.5), Inches(y), Inches(9), Inches(0.4), size=11, bold=True)
    y += 0.4
    for item in subitems:
        add_text(sl, '  ' + item, Inches(0.5), Inches(y), Inches(9), Inches(0.38), size=10)
        y += 0.38
    y += 0.1

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Zakljucak
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Zakljucak')

conclusions = [
    'QoS 0 daje 7x veci throughput od QoS 2 (623k vs 43k msg/s) uz zanemarljiv CPU',
    'QoS 2 trosi ceo CPU core na 10 000 uredjaja (102%) zbog dvostrukog handshake-a',
    'p50 latencija je gotovo ista za sva 3 QoS nivoa (~4ms) — razlika je na repovima (p99)',
    'Mosquitto broker nije usko grlo — Spring klijent je 750x sporiji od direktnog benchmarka',
    'Burst load: sistem se odmah adaptira, CPU se vraca na bazni nivo bez akumulacije',
    'cleanSession=false + automaticReconnect = pouzdana isporuka pri padu mreze',
    'MQTT je optimalan za edge ingestion, Kafka za cloud storage i stream processing',
]

for i, c in enumerate(conclusions):
    add_text(sl, f'{i+1}.  {c}',
             Inches(0.5), Inches(1.1 + i*0.57), Inches(9), Inches(0.5),
             size=10)

out_path = os.path.join(OUT, 'prezentacija.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
