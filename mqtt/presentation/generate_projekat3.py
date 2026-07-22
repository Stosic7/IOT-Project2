"""
Generator prezentacije za Projekat 3 (eKuiper CEP + MaaS).

Pokretanje (sistemski python3 — ima matplotlib + python-pptx + numpy):
    cd mqtt/presentation && python3 generate_projekat3.py

Napomena: namerno NE koristi pandas (nije instaliran u sistemskom python-u) —
CSV se cita `csv` modulom iz stdlib-a.
"""
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import csv
import io
import os

OUT = os.path.dirname(os.path.abspath(__file__))
CSV_PATH = os.path.join(OUT, '..', 'src', 'main', 'resources', 'real_time_data.csv')
DASHBOARD_PNG = os.path.join(OUT, 'dashboard.png')
MAX_ROWS = 150_000            # isti uzorak koji koristi i train.py

# ── Tema (identicna Projektu 2) ─────────────────────────────────────────
BLACK  = RGBColor(0x00, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MUTED  = RGBColor(0x66, 0x66, 0x66)
ACCENT = RGBColor(0x1A, 0x56, 0xDB)
GREEN  = RGBColor(0x0B, 0x7A, 0x3B)
RED    = RGBColor(0xB3, 0x1D, 0x1D)

CHART_COLORS = ['#1A56DB', '#4B9CD3', '#A8C8F0']
C_HOT, C_COLD = '#C0392B', '#1A56DB'

# ── Izmereni podaci ─────────────────────────────────────────────────────
# ML metrike — tacan ispis `python train.py` (hronoloski split 70/15/15, 150k redova)
ML = {
    'LinearRegression\n(baseline)': {'val_mae': 2.002, 'val_r2': 0.488, 'mae': 2.001, 'r2': 0.504},
    'RandomForest\n(deploy-ovan)':  {'val_mae': 2.004, 'val_r2': 0.486, 'mae': 2.004, 'r2': 0.503},
}

# Live merenje (100 uredjaja, QoS 1) — isti run sa kog je i screenshot dashboarda
LIVE = {
    'sent': 971_600, 'db_records': 3_865_809, 'batch': 500,
    'events': 1_998, 'ev_high': 1_066, 'ev_low': 932,
    'window_msgs': 99_416, 'avg_temp': 20.08, 'predicted': 22.00, 'delta': 1.92,
    'p50': 3, 'p95': 54, 'p99': 99,
}

TH_HIGH, TH_LOW = 29.5, 8.5


# ── Helperi ─────────────────────────────────────────────────────────────
def set_bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_title(slide, text, top=Inches(0.32), size=26):
    txb = slide.shapes.add_textbox(Inches(0.5), top, Inches(9.2), Inches(0.6))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = BLACK


def add_text(slide, text, left, top, width, height, size=11, bold=False,
             align=PP_ALIGN.LEFT, color=BLACK, italic=False, font=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if font:
        run.font.name = font
    return txb


def add_bullets(slide, items, left=0.55, top=1.05, width=9.0, size=11, gap=0.42):
    """items: lista strin­gova ili (tekst, bold, color) tuple-ova."""
    y = top
    for it in items:
        if isinstance(it, tuple):
            text, bold, color = it
        else:
            text, bold, color = it, False, BLACK
        if text == '':
            y += gap * 0.45
            continue
        add_text(slide, text, Inches(left), Inches(y), Inches(width), Inches(gap),
                 size=size, bold=bold, color=color)
        y += gap
    return y


def add_conclusion(slide, text, top=4.82):
    """Plava traka sa zakljuckom na dnu slajda."""
    bar = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9.0), Inches(0.5))
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = 'Zaključak:  '
    r1.font.size = Pt(10.5)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(10.5)
    r2.font.color.rgb = BLACK


def chart_to_image(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf


def style_ax(ax, title=''):
    ax.set_facecolor('white')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#CCCCCC')
    ax.spines['bottom'].set_color('#CCCCCC')
    ax.tick_params(colors='#333333', labelsize=9)
    if title:
        ax.set_title(title, fontsize=10, fontweight='bold', color='#111111', pad=8)


def load_temperatures(path, limit=MAX_ROWS):
    """Cita kolonu temperature iz CSV-a (stdlib csv — bez pandas-a)."""
    temps = []
    with open(path, newline='', encoding='utf-8') as f:
        reader = csv.reader(f)
        next(reader, None)  # header
        for i, row in enumerate(reader):
            if i >= limit:
                break
            try:
                temps.append(float(row[2]))
            except (IndexError, ValueError):
                continue
    return np.array(temps)


# ── Podaci iz dataset-a ─────────────────────────────────────────────────
temps = load_temperatures(CSV_PATH)
T_AVG, T_MIN, T_MAX = temps.mean(), temps.min(), temps.max()
PCT_HIGH = (temps > TH_HIGH).mean() * 100
PCT_LOW = (temps < TH_LOW).mean() * 100
PCT_EVENTS = PCT_HIGH + PCT_LOW
print(f'Dataset: n={len(temps):,}  avg={T_AVG:.2f}  min={T_MIN:.2f}  max={T_MAX:.2f}')
print(f'  >{TH_HIGH}°C = {PCT_HIGH:.2f}%   <{TH_LOW}°C = {PCT_LOW:.2f}%')

# ── Prezentacija ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]

# ═══════════════════════════════════════════════ 1 — Naslovna
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_text(sl, 'Projekat 3', Inches(1), Inches(1.35), Inches(8), Inches(0.4),
         size=13, bold=True, align=PP_ALIGN.CENTER, color=ACCENT)
add_text(sl, 'Unapređenje Analytics mikroservisa:\neKuiper CEP + MaaS',
         Inches(1), Inches(1.8), Inches(8), Inches(1.1),
         size=28, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, 'Streaming obrada događaja i mašinsko učenje kao servis nad MQTT tokom podataka',
         Inches(1), Inches(3.05), Inches(8), Inches(0.5),
         size=13, align=PP_ALIGN.CENTER, color=RGBColor(0x44, 0x44, 0x44))
add_text(sl, 'Spring Boot  ·  eKuiper  ·  FastAPI + scikit-learn  ·  React  ·  Docker Compose',
         Inches(1), Inches(3.7), Inches(8), Inches(0.4),
         size=11, align=PP_ALIGN.CENTER, color=MUTED)

# ═══════════════════════════════════════════════ 2 — Zahtevi → status
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Zahtevi zadatka i šta je implementirano')

rows = [
    ('1a. eKuiper CEP preko MQTT brokera',
     'Stream nad iot/sensors/# + 2 SQL pravila, sink na iot/events'),
    ('1b. MaaS i njegovi REST endpointi',
     'Analytics zove POST /predict na kraju svakog prozora'),
    ('2. eKuiper na istom topicu, šalje na novi',
     'iot/sensors/#  →  detekcija  →  iot/events  →  Analytics'),
    ('3. Python + FastAPI + scikit-learn model',
     'RandomForestRegressor, trening/validacija/test 70/15/15'),
    ('4. Docker kontejneri + Web aplikacija',
     '7 kontejnera; React dashboard sa live prikazom'),
    ('5. Izvorni kod na GitHub-u',
     'Commit-ovano i push-ovano na origin/main'),
]
col_x, col_w = [0.55, 4.75], [4.1, 4.7]
y = 1.05
add_text(sl, 'Zahtev', Inches(col_x[0]), Inches(y), Inches(col_w[0]), Inches(0.3), size=10, bold=True, color=MUTED)
add_text(sl, 'Implementacija', Inches(col_x[1]), Inches(y), Inches(col_w[1]), Inches(0.3), size=10, bold=True, color=MUTED)
y += 0.38
for req, impl in rows:
    add_text(sl, req, Inches(col_x[0]), Inches(y), Inches(col_w[0]), Inches(0.5), size=10, bold=True)
    add_text(sl, impl, Inches(col_x[1]), Inches(y), Inches(col_w[1]), Inches(0.5), size=10, color=RGBColor(0x33, 0x33, 0x33))
    y += 0.58

add_conclusion(sl, 'svih 5 tačaka je pokriveno; MaaS koristi FastAPI, što zadatak izričito dozvoljava.')

# ═══════════════════════════════════════════════ 3 — Arhitektura
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Arhitektura — šta je dodato u odnosu na Projekat 2')

fig, ax = plt.subplots(figsize=(10, 3.5))
fig.patch.set_facecolor('white')
ax.set_xlim(0, 100)
ax.set_ylim(0, 42)
ax.axis('off')


def box(x, y, w, h, label, sub='', color='#1A56DB', fill='#FFFFFF', lw=1.6, fs=9):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.6",
                                linewidth=lw, edgecolor=color, facecolor=fill))
    ax.text(x + w / 2, y + h / 2 + (1.6 if sub else 0), label,
            ha='center', va='center', fontsize=fs, fontweight='bold', color='#111111')
    if sub:
        ax.text(x + w / 2, y + h / 2 - 2.6, sub, ha='center', va='center',
                fontsize=7.5, color='#555555')


def arrow(x1, y1, x2, y2, label='', color='#666666', style='-|>', ls='-'):
    ax.add_patch(FancyArrowPatch((x1, y1), (x2, y2), arrowstyle=style,
                                 mutation_scale=12, linewidth=1.3, color=color, linestyle=ls))
    if label:
        ax.text((x1 + x2) / 2, (y1 + y2) / 2 + 1.6, label, ha='center', va='bottom',
                fontsize=7, color=color, style='italic')


# Postojece (Projekat 2) — sivo
box(1, 26, 15, 9, 'Ingestion', '100 uređaja', color='#999999')
box(21, 26, 15, 9, 'Mosquitto', 'iot/sensors/#', color='#999999')
box(43, 34, 16, 7.5, 'Storage', 'batch 500', color='#999999')
box(64, 34, 15, 7.5, 'PostgreSQL', '', color='#999999')
box(43, 22, 16, 8.5, 'Analytics', 'tumbling 10s', color='#999999')

# Novo (Projekat 3) — plavo
box(43, 8, 16, 8.5, 'eKuiper', 'CEP pravila', color='#1A56DB', fill='#EEF4FE', lw=2)
box(70, 22, 15, 8.5, 'MaaS', 'FastAPI /predict', color='#1A56DB', fill='#EEF4FE', lw=2)

arrow(16.5, 30.5, 20.5, 30.5)
arrow(36.5, 32, 42.5, 37)          # -> Storage
arrow(36.5, 30.5, 42.5, 27)        # -> Analytics
arrow(36.5, 29, 42.5, 13, color='#1A56DB')   # -> eKuiper
arrow(59.5, 37.5, 63.5, 37.5)      # Storage -> PG
arrow(51, 17, 51, 21.5, color='#1A56DB')  # eKuiper -> Analytics (preko brokera)
ax.text(52.4, 19.0, 'iot/events', ha='left', va='center', fontsize=7,
        color='#1A56DB', style='italic')
arrow(59.5, 26.5, 69.5, 26.5, label='REST', color='#1A56DB', style='<|-|>')

ax.text(1, 4, 'sivo = postojalo u Projektu 2       plavo = dodato u Projektu 3',
        fontsize=8, color='#555555')
plt.tight_layout(pad=0.3)
sl.shapes.add_picture(chart_to_image(fig), Inches(0.25), Inches(0.95), Inches(9.5), Inches(3.6))

add_conclusion(sl, 'Projekat 3 je dodao samo 2 kontejnera i proširio Analytics — Ingestion, '
                   'Storage i baza rade nepromenjeni.')

# ═══════════════════════════════════════════════ 4 — eKuiper bez koda
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'eKuiper — CEP bez ijedne linije aplikativnog koda')

add_bullets(sl, [
    ('Sve se konfiguriše kroz REST API (POST /streams, POST /rules) — nula Java/Python koda:', True, BLACK),
], top=1.02, size=11)

add_text(sl, 'CREATE STREAM sensors () WITH (DATASOURCE="iot/sensors/#",\n'
             '                              FORMAT="JSON", TYPE="mqtt")',
         Inches(0.7), Inches(1.5), Inches(8.6), Inches(0.7), size=9.5, font='Consolas',
         color=RGBColor(0x1A, 0x3A, 0x6B))

add_text(sl, 'SELECT device_id, temperature, location, "HIGH_TEMPERATURE" AS event_type\n'
             'FROM sensors WHERE temperature > 29.5      →  sink: mqtt topic "iot/events"',
         Inches(0.7), Inches(2.32), Inches(8.6), Inches(0.7), size=9.5, font='Consolas',
         color=RGBColor(0x1A, 0x3A, 0x6B))

add_bullets(sl, [
    ('Ista pravila važe za svaku poruku, bez stanja — filtriranje po poruci (stateless).', False, BLACK),
    ('Analytics ne zna da eKuiper postoji — komunikacija ide isključivo preko MQTT topica.', False, BLACK),
    ('Pravilo se dodaje/menja jednim curl pozivom, dok sistem radi.', False, BLACK),
], top=3.25, size=11, gap=0.42)

add_conclusion(sl, 'pravila se menjaju u runtime-u bez ponovnog build-a i deploy-a Analytics servisa.')

# ═══════════════════════════════════════════════ 5 — Kalibracija pragova
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Kalibracija pragova — pragovi se izvode iz podataka')

fig, ax = plt.subplots(figsize=(10, 3.1))
fig.patch.set_facecolor('white')
n, bins, patches = ax.hist(temps, bins=90, color='#B9CFF0', edgecolor='white', linewidth=0.4)
for patch, left_edge in zip(patches, bins[:-1]):
    if left_edge >= TH_HIGH:
        patch.set_facecolor(C_HOT)
    elif left_edge < TH_LOW:
        patch.set_facecolor(C_COLD)
ax.axvline(TH_HIGH, color=C_HOT, linestyle='--', linewidth=1.4)
ax.axvline(TH_LOW, color=C_COLD, linestyle='--', linewidth=1.4)
ax.text(TH_HIGH + 0.4, ax.get_ylim()[1] * 0.82, f'> {TH_HIGH}°C\n{PCT_HIGH:.2f}%',
        fontsize=8.5, color=C_HOT, fontweight='bold')
ax.text(TH_LOW - 0.4, ax.get_ylim()[1] * 0.82, f'< {TH_LOW}°C\n{PCT_LOW:.2f}%',
        fontsize=8.5, color=C_COLD, fontweight='bold', ha='right')
ax.set_xlabel('Temperatura (°C)', fontsize=9)
ax.set_ylabel('Broj očitavanja', fontsize=9)
style_ax(ax, f'Raspodela temperature u dataset-u (n = {len(temps):,};  '
             f'prosek {T_AVG:.2f}°C, min {T_MIN:.2f}°C, max {T_MAX:.2f}°C)')
plt.tight_layout(pad=0.4)
sl.shapes.add_picture(chart_to_image(fig), Inches(0.3), Inches(0.95), Inches(9.4), Inches(2.95))

add_bullets(sl, [
    (f'Pragovi 29.5°C / 8.5°C biraju samo repove raspodele — ukupno {PCT_EVENTS:.2f}% poruka.', False, BLACK),
    ('Prvobitno pravilo LOW_BATTERY je odbačeno: baterija je u dataset-u uvek 80–100%, '
     'pa ne bi okinulo nijednom.', False, BLACK),
], top=4.02, size=10.5, gap=0.36)

add_conclusion(sl, 'prag koji se „pogodi napamet” ili nikad ne okine ili preplavi sistem — '
                   'mora se izvesti iz raspodele podataka.')

# ═══════════════════════════════════════════════ 6 — CEP u radu
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'CEP u radu — filter koji rasterećuje sistem')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.2))
fig.patch.set_facecolor('white')

bars = axes[0].bar(['HIGH_TEMPERATURE', 'LOW_TEMPERATURE'],
                   [LIVE['ev_high'], LIVE['ev_low']],
                   color=[C_HOT, C_COLD], width=0.5)
for b, v in zip(bars, [LIVE['ev_high'], LIVE['ev_low']]):
    axes[0].text(b.get_x() + b.get_width() / 2, b.get_height() + 25, f'{v:,}',
                 ha='center', fontsize=9, fontweight='bold')
axes[0].set_ylabel('Broj događaja')
axes[0].set_ylim(0, max(LIVE['ev_high'], LIVE['ev_low']) * 1.25)
style_ax(axes[0], f'Detektovani događaji (ukupno {LIVE["events"]:,})')

share = LIVE['events'] / LIVE['sent'] * 100
axes[1].barh(['Poruke\nna brokeru', 'Događaji\n(iot/events)'],
             [LIVE['sent'], LIVE['events']], color=['#B9CFF0', '#1A56DB'], height=0.5)
axes[1].set_xscale('log')
axes[1].set_xlabel('broj poruka (logaritamska skala)')
axes[1].text(LIVE['sent'], 0, f'  {LIVE["sent"]:,}', va='center', fontsize=9)
axes[1].text(LIVE['events'], 1, f'  {LIVE["events"]:,}  ({share:.2f}%)', va='center',
             fontsize=9, fontweight='bold', color='#1A56DB')
axes[1].set_xlim(1, LIVE['sent'] * 12)
style_ax(axes[1], 'Udeo događaja u ukupnom saobraćaju')

plt.tight_layout(pad=0.8)
sl.shapes.add_picture(chart_to_image(fig), Inches(0.3), Inches(0.98), Inches(9.4), Inches(3.05))

add_bullets(sl, [
    (f'Od {LIVE["sent"]:,} poruka, eKuiper je propustio samo {LIVE["events"]:,} '
     f'događaja od interesa ({share:.2f}%).', False, BLACK),
], top=4.18, size=10.5)

add_conclusion(sl, 'CEP sloj smanjuje saobraćaj ka potrošačima ~500×, pa Analytics obrađuje '
                   'događaje umesto sirovih očitavanja.')

# ═══════════════════════════════════════════════ 7 — Neuspeli prvi pristup
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'MaaS — prvi pristup modelu nije uspeo')

add_bullets(sl, [
    ('Prvobitna ideja: predvideti sledeću temperaturu iz prethodnih 10 očitavanja uređaja', True, RED),
    ('Rezultat: negativan R² — model je bio gori od običnog proseka.', False, BLACK),
    ('Razlog: temperatura po uređaju se ponaša kao šum — susedna očitavanja skaču 25°C → 18°C,', False, BLACK),
    ('pa lag-featuri (istorija) ne nose nikakvu informaciju.', False, BLACK),
    '',
    ('Pivot: regresija iz ostalih senzorskih veličina u istoj poruci („virtuelni senzor”)', True, GREEN),
    ('Signal postoji između veličina, a ne kroz vreme:', False, BLACK),
    ('     ·  location — Outside ≈ 15°C  vs  sobe ≈ 22°C  (najjači prediktor)', False, BLACK),
    ('     ·  humidity  r = −0,27          ·  light  r = −0,21', False, BLACK),
], top=1.05, size=11, gap=0.40)

add_conclusion(sl, 'to što su podaci vremenska serija ne znači da su vremenski predvidivi — '
                   'signal je bio između senzora, a ne kroz vreme.')

# ═══════════════════════════════════════════════ 8 — Model i metrike
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'MaaS — model, trening i rezultati')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.0))
fig.patch.set_facecolor('white')
names = list(ML.keys())
x = np.arange(len(names))
w = 0.34

axes[0].bar(x - w / 2, [ML[n]['val_mae'] for n in names], w, label='validacija', color=CHART_COLORS[1])
axes[0].bar(x + w / 2, [ML[n]['mae'] for n in names], w, label='test', color=CHART_COLORS[0])
for i, n in enumerate(names):
    axes[0].text(i - w / 2, ML[n]['val_mae'] + 0.03, f"{ML[n]['val_mae']:.3f}", ha='center', fontsize=8)
    axes[0].text(i + w / 2, ML[n]['mae'] + 0.03, f"{ML[n]['mae']:.3f}", ha='center', fontsize=8, fontweight='bold')
axes[0].set_xticks(x); axes[0].set_xticklabels(names, fontsize=8.5)
axes[0].set_ylabel('MAE (°C)'); axes[0].set_ylim(0, 2.6)
axes[0].legend(fontsize=8)
style_ax(axes[0], 'Srednja apsolutna greška — manje je bolje')

axes[1].bar(x - w / 2, [ML[n]['val_r2'] for n in names], w, label='validacija', color=CHART_COLORS[1])
axes[1].bar(x + w / 2, [ML[n]['r2'] for n in names], w, label='test', color=CHART_COLORS[0])
for i, n in enumerate(names):
    axes[1].text(i - w / 2, ML[n]['val_r2'] + 0.012, f"{ML[n]['val_r2']:.3f}", ha='center', fontsize=8)
    axes[1].text(i + w / 2, ML[n]['r2'] + 0.012, f"{ML[n]['r2']:.3f}", ha='center', fontsize=8, fontweight='bold')
axes[1].set_xticks(x); axes[1].set_xticklabels(names, fontsize=8.5)
axes[1].set_ylabel('R²'); axes[1].set_ylim(0, 0.68)
axes[1].legend(fontsize=8)
style_ax(axes[1], 'Koeficijent determinacije — više je bolje')

plt.tight_layout(pad=0.8)
sl.shapes.add_picture(chart_to_image(fig), Inches(0.3), Inches(0.95), Inches(9.4), Inches(2.9))

add_bullets(sl, [
    ('150 000 redova  ·  9 featura  ·  hronološki split 70/15/15 (bez mešanja — vremenska serija)', False, BLACK),
    ('RandomForest (MAE 2,004) nije nadmašio linearni baseline (MAE 2,001) — razlika je u trećoj decimali.', False, BLACK),
], top=4.02, size=10.5, gap=0.36)

add_conclusion(sl, 'složeniji model nije doneo ništa — granicu je postavio kvalitet featura, '
                   'a ne izbor algoritma.')

# ═══════════════════════════════════════════════ 9 — Integracija MaaS-a
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Integracija MaaS-a — jedan poziv po prozoru, ne po poruci')

add_bullets(sl, [
    ('Ključna projektna odluka:', True, BLACK),
    (f'Analytics agregira ceo prozor (~{LIVE["window_msgs"]:,} poruka) i tek onda šalje '
     f'JEDAN REST poziv.', False, BLACK),
    ('Poziv po poruci bi značio ~100 000 HTTP zahteva na svakih 10 sekundi — neizvodljivo.', False, BLACK),
    '',
    ('Šta se šalje:', True, BLACK),
    ('Proseci humidity / pressure / light / sound / motion + najčešća lokacija iz prozora.', False, BLACK),
    '',
    ('Otpornost na otkaz:', True, BLACK),
    ('Poziv je u try/catch — ako MaaS padne, Analytics nastavlja da radi bez predikcije.', False, BLACK),
    ('MaaS je bez stanja (učitan model.joblib), pa se skalira i restartuje nezavisno.', False, BLACK),
], top=1.02, size=11, gap=0.40)

add_conclusion(sl, 'agregacioni prozor je prirodna tačka integracije — spaja tempo streaminga '
                   'sa tempom sinhronog REST modela.')

# ═══════════════════════════════════════════════ 10 — Predvidjeno vs stvarno
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Predviđeno vs stvarno — validacija u radu')

fig, ax = plt.subplots(figsize=(9.2, 2.9))
fig.patch.set_facecolor('white')
labels = ['Predviđeno\n(MaaS model)', 'Stvarni prosek\n(tumbling window)']
vals = [LIVE['predicted'], LIVE['avg_temp']]
bars = ax.bar(labels, vals, color=['#7C6CF0', '#1A56DB'], width=0.42)
for b, v in zip(bars, vals):
    ax.text(b.get_x() + b.get_width() / 2, v + 0.35, f'{v:.2f} °C',
            ha='center', fontsize=11, fontweight='bold')
ax.set_ylabel('Temperatura (°C)')
ax.set_ylim(0, 26)
ax.set_xlim(-0.6, 1.75)
# osencen pojas izmedju predvidjene i stvarne vrednosti = greska modela
ax.axhspan(LIVE['avg_temp'], LIVE['predicted'], color='#C0392B', alpha=0.11, zorder=0)
ax.hlines([LIVE['avg_temp'], LIVE['predicted']], -0.45, 1.72,
          colors='#C0392B', linestyles='--', linewidth=1, zorder=1)
ax.text(1.70, (LIVE['predicted'] + LIVE['avg_temp']) / 2,
        f'Δ {LIVE["delta"]:.2f} °C', ha='right', va='center', fontsize=10.5,
        fontweight='bold', color='#C0392B')
style_ax(ax, 'Merenje uživo na kraju jednog 10-sekundnog prozora')
plt.tight_layout(pad=0.4)
sl.shapes.add_picture(chart_to_image(fig), Inches(0.75), Inches(0.95), Inches(8.5), Inches(2.8))

rf_mae = ML['RandomForest\n(deploy-ovan)']['mae']
add_bullets(sl, [
    (f'Odstupanje u radu Δ {LIVE["delta"]:.2f} °C  ≈  MAE {rf_mae:.3f} °C '
     f'izmeren offline na test skupu.', True, BLACK),
    ('Model se u produkciji ponaša isto kao na test skupu — nema drift-a ni greške u integraciji.', False, BLACK),
], top=3.92, size=10.5, gap=0.38)

add_conclusion(sl, 'poklapanje live odstupanja sa offline MAE potvrđuje da je cela putanja '
                   'podataka do modela ispravna.')

# ═══════════════════════════════════════════════ 11 — Performanse
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Performanse — CEP i MaaS ne usporavaju glavni tok')

fig, axes = plt.subplots(1, 2, figsize=(10, 3.0))
fig.patch.set_facecolor('white')

lat = [LIVE['p50'], LIVE['p95'], LIVE['p99']]
bars = axes[0].bar(['p50', 'p95', 'p99'], lat, color=CHART_COLORS, width=0.5)
for b, v in zip(bars, lat):
    axes[0].text(b.get_x() + b.get_width() / 2, v + 2, f'{v} ms', ha='center',
                 fontsize=9.5, fontweight='bold')
axes[0].set_ylabel('ms')
axes[0].set_ylim(0, max(lat) * 1.3)
style_ax(axes[0], 'End-to-end latencija (slanje → obrada)')

names2 = ['Poslato\n(ingestion)', 'Upisano\n(PostgreSQL)', 'Poruka\nu prozoru']
vals2 = [LIVE['sent'], LIVE['db_records'], LIVE['window_msgs']]
bars = axes[1].bar(names2, vals2, color=['#B9CFF0', '#4B9CD3', '#1A56DB'], width=0.5)
for b, v in zip(bars, vals2):
    axes[1].text(b.get_x() + b.get_width() / 2, v * 1.05, f'{v:,}', ha='center', fontsize=8.5)
axes[1].set_yscale('log')
axes[1].set_ylabel('broj poruka (log)')
style_ax(axes[1], 'Obim obrađenih podataka')

plt.tight_layout(pad=0.8)
sl.shapes.add_picture(chart_to_image(fig), Inches(0.3), Inches(0.95), Inches(9.4), Inches(2.9))

add_bullets(sl, [
    ('eKuiper obrađuje isti tok paralelno, a MaaS se zove jednom u 10 s — nijedan nije '
     'na putanji poruke.', False, BLACK),
    (f'p50 ostaje {LIVE["p50"]} ms i pri ~{LIVE["window_msgs"]:,} poruka po prozoru; '
     f'rep (p99 {LIVE["p99"]} ms) dolazi od batch upisa u bazu.', False, BLACK),
], top=4.02, size=10.5, gap=0.36)

add_conclusion(sl, 'dodavanje CEP-a i ML servisa nije pomerilo medijalnu latenciju — '
                   'oba rade van hot path-a.')

# ═══════════════════════════════════════════════ 12 — Dashboard
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Web aplikacija — sve tri komponente u jednom pogledu')

if os.path.exists(DASHBOARD_PNG):
    sl.shapes.add_picture(DASHBOARD_PNG, Inches(0.45), Inches(1.0), height=Inches(4.25))

add_bullets(sl, [
    ('React + Vite, polling na 2 s', True, BLACK),
    ('Arhitektura toka — živi brojači po', False, BLACK),
    ('    komponenti (događaji, poruke, °C)', False, BLACK),
    ('MaaS kartica — predviđeno vs stvarno', False, BLACK),
    ('    sa odstupanjem i pragom alarma', False, BLACK),
    ('eKuiper kartica — ukupno događaja,', False, BLACK),
    ('    podela HIGH/LOW i live feed', False, BLACK),
    ('Tumbling window — prosek, alarm', False, BLACK),
    ('    i p50/p95/p99 latencija', False, BLACK),
    '',
    ('Demo bez backend-a: ?mock=1', True, ACCENT),
], left=4.9, top=1.15, width=4.7, size=10, gap=0.32)

add_conclusion(sl, 'dashboard pokriva zahtev za Web aplikacijom i služi kao alat za '
                   'demonstraciju sistema uživo.')

# ═══════════════════════════════════════════════ 13 — Zakljucci
sl = prs.slides.add_slide(blank)
set_bg_white(sl)
add_title(sl, 'Zaključci')

conclusions = [
    'eKuiper daje CEP bez ijedne linije koda — SQL pravila preko REST API-ja, izmenjiva u runtime-u.',
    'MQTT topic kao granica servisa: Analytics i eKuiper se ne poznaju, vezuje ih samo iot/events.',
    'Pragovi pravila moraju se izvesti iz raspodele podataka — inače nikad ne okinu ili preplave sistem.',
    'CEP je smanjio saobraćaj ka potrošačima ~500× (0,21% poruka je postalo događaj).',
    'Vremenska serija nije automatski predvidiva — lag-featuri su dali negativan R², signal je bio '
    'između senzora (location, humidity, light).',
    'Složeniji model nije pomogao: RandomForest ≈ linearni baseline (MAE 2,004 vs 2,001).',
    'Live odstupanje Δ 1,92 °C poklopilo se sa offline MAE ≈ 2,0 °C — integracija je ispravna.',
    'MaaS se zove jednom po prozoru i u try/catch bloku — sistem radi i kada model nije dostupan.',
]
y = 1.02
for i, c in enumerate(conclusions):
    add_text(sl, f'{i + 1}.', Inches(0.55), Inches(y), Inches(0.35), Inches(0.5),
             size=10.5, bold=True, color=ACCENT)
    add_text(sl, c, Inches(0.95), Inches(y), Inches(8.5), Inches(0.5), size=10.5)
    y += 0.55 if len(c) < 105 else 0.62

out_path = os.path.join(OUT, 'prezentacija_projekat3.pptx')
prs.save(out_path)
print(f'Sačuvano: {out_path}  ({len(prs.slides._sldIdLst)} slajdova)')
